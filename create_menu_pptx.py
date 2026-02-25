#!/usr/bin/env python3
"""Dify AI ワークフロー ファミレスメニュー 提案資料 PowerPoint生成スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os

# ── カラーパレット ──
NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
BLUE      = RGBColor(0x2D, 0x6C, 0xDF)
LIGHT_BLUE= RGBColor(0xE8, 0xF0, 0xFE)
ACCENT    = RGBColor(0xFF, 0x6B, 0x35)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY  = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY= RGBColor(0xF5, 0xF5, 0xF5)
GREEN     = RGBColor(0x27, 0xAE, 0x60)
RED_ACCENT= RGBColor(0xE7, 0x4C, 0x3C)
GOLD      = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── ヘルパー関数 ──
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=14,
                          color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                          line_spacing=1.5, font_name="Meiryo"):
    """lines: list of (text, font_size, color, bold) tuples or just strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.size = Pt(line[1]) if len(line) > 1 else Pt(font_size)
            p.font.color.rgb = line[2] if len(line) > 2 else color
            p.font.bold = line[3] if len(line) > 3 else bold
        else:
            p.text = str(line)
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox

def add_circle_number(slide, left, top, size, number, bg_color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(int(size / Pt(1) * 0.45))
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Meiryo"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def slide_header(slide, title, subtitle=None, accent_bar=True):
    """共通ヘッダー"""
    # Top bar
    add_shape(slide, Inches(0), Inches(0), W, Inches(0.06), BLUE)
    # Title
    add_textbox(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                title, font_size=30, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.4),
                    subtitle, font_size=16, color=MID_GRAY)
    # Accent bar under title
    if accent_bar:
        add_shape(slide, Inches(0.8), Inches(1.25), Inches(1.5), Inches(0.04), ACCENT)

def slide_footer(slide, page_num):
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(5), Inches(0.3),
                "Dify AI ワークフロー ファミレスメニュー  |  Confidential",
                font_size=9, color=MID_GRAY)
    add_textbox(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
                str(page_num), font_size=9, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


def make_table(slide, left, top, width, rows_data, col_widths, header_color=NAVY,
               font_size=11, row_height=0.42):
    """rows_data: list of lists (first row = header)"""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_height = Inches(row_height * n_rows)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_height)
    table = table_shape.table

    # Column widths
    for i, cw in enumerate(col_widths):
        table.columns[i].width = Inches(cw)

    for r_idx, row_data in enumerate(rows_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Meiryo"
                if r_idx == 0:
                    paragraph.font.color.rgb = WHITE
                    paragraph.font.bold = True
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_GRAY
                    if c_idx == 0:
                        paragraph.alignment = PP_ALIGN.CENTER
                    else:
                        paragraph.alignment = PP_ALIGN.LEFT

            # Cell fill
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


# ════════════════════════════════════════════════════
# SLIDE 1: 表紙
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, NAVY)

# Accent bar
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
add_shape(slide, Inches(0), Inches(7.42), W, Inches(0.08), ACCENT)

# Title area
add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(1.0),
            "Dify AI ワークフロー", font_size=24, color=RGBColor(0xAA,0xCC,0xFF), bold=False)
add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
            "ファミレスメニュー", font_size=52, color=WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.7),
            "〜 業界別 AI ソリューション カタログ 全170品 〜",
            font_size=22, color=ACCENT, bold=True)

# Separator
add_shape(slide, Inches(1.5), Inches(4.0), Inches(3), Inches(0.03), ACCENT)

# Subtitle lines
add_multiline_textbox(slide, Inches(1.5), Inches(4.3), Inches(9), Inches(1.5), [
    ("20業界 × 170ワークフロー の網羅的AIソリューションで", 18, RGBColor(0xCC,0xDD,0xFF), False),
    ("御社の業務効率化・コスト削減・売上拡大を実現します", 18, RGBColor(0xCC,0xDD,0xFF), False),
], line_spacing=1.8)

# Date & Company
today = datetime.date.today().strftime("%Y年%m月")
add_textbox(slide, Inches(1.5), Inches(6.2), Inches(5), Inches(0.4),
            f"{today}  |  Powered by Dify", font_size=14, color=MID_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 2: 企業のAI導入における課題
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "企業のAI活用における3つの壁", "なぜ多くの企業でAI導入が「PoC止まり」になるのか？")
slide_footer(slide, 2)

# 3 columns for pain points
pains = [
    ("1", "何を作ればいいかわからない",
     "AI活用のアイデアはあるが、\n自社の業務にどう適用すれば\n効果があるのか具体像が見えない",
     "「AIで何かできそう」\n→ 具体化できず停滞"),
    ("2", "作っても使われない",
     "技術先行で構築したAIツールが\n現場の業務フローに合わず、\n利用率が上がらない",
     "「便利そうだけど今のやり方で\n  十分」→ 現場が使わない"),
    ("3", "コストが見合わない",
     "個別開発すると1機能あたり\n数百万円〜。ROIが見えず\n経営層の承認が得られない",
     "「費用対効果を示せ」\n→ 予算がつかない"),
]

for i, (num, title, desc, quote) in enumerate(pains):
    x = Inches(0.8 + i * 4.0)
    # Card background
    card = add_rounded_rect(slide, x, Inches(1.6), Inches(3.6), Inches(5.0), LIGHT_GRAY)

    # Number circle
    add_circle_number(slide, x + Inches(1.4), Inches(1.85), Inches(0.7), num, RED_ACCENT)

    # Pain title
    add_textbox(slide, x + Inches(0.2), Inches(2.7), Inches(3.2), Inches(0.7),
                title, font_size=20, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_multiline_textbox(slide, x + Inches(0.3), Inches(3.4), Inches(3.0), Inches(1.8),
                          [desc], font_size=13, color=DARK_GRAY, line_spacing=1.6)

    # Quote box
    qbox = add_rounded_rect(slide, x + Inches(0.2), Inches(5.2), Inches(3.2), Inches(1.1),
                             RGBColor(0xFF,0xEB,0xEE))
    add_multiline_textbox(slide, x + Inches(0.4), Inches(5.3), Inches(2.8), Inches(0.9),
                          [quote], font_size=11, color=RED_ACCENT, line_spacing=1.4)


# ════════════════════════════════════════════════════
# SLIDE 3: Difyとファミレスメニューで解決
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "解決策：Dify × ファミレスメニュー",
             "「何を作るか」は、もう考えなくていい")
slide_footer(slide, 3)

# Left: Dify explanation
add_rounded_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5), LIGHT_BLUE)
add_textbox(slide, Inches(1.0), Inches(1.7), Inches(5.0), Inches(0.5),
            "Dify とは", font_size=22, color=BLUE, bold=True)
add_multiline_textbox(slide, Inches(1.0), Inches(2.2), Inches(5.0), Inches(1.6), [
    ("ノーコード / ローコードで AI ワークフローを構築・運用できる", 14, DARK_GRAY, False),
    ("オープンソースの AI アプリケーション開発プラットフォーム。", 14, DARK_GRAY, False),
    ("", 8, DARK_GRAY, False),
    ("Claude / GPT / Gemini 等の LLM を組み合わせ、", 14, DARK_GRAY, False),
    ("RAG・分岐・API連携・バッチ処理を GUI で構築可能。", 14, DARK_GRAY, False),
], line_spacing=1.5)

# Right: Menu explanation
add_rounded_rect(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5), RGBColor(0xFF,0xF3,0xE0))
add_textbox(slide, Inches(7.2), Inches(1.7), Inches(5.0), Inches(0.5),
            "ファミレスメニュー とは", font_size=22, color=ACCENT, bold=True)
add_multiline_textbox(slide, Inches(7.2), Inches(2.2), Inches(5.0), Inches(1.6), [
    ("170品の業界別AIワークフローを「メニュー表」として用意。", 14, DARK_GRAY, False),
    ("お客様は「メニューから選ぶだけ」で、", 14, DARK_GRAY, False),
    ("", 8, DARK_GRAY, False),
    ("自社に最適なAIワークフローを", 14, DARK_GRAY, True),
    ("最短1日で導入開始 できます。", 14, ACCENT, True),
], line_spacing=1.5)

# Bottom: 3 value cards
values = [
    ("選ぶだけで導入", "170品のメニューから\n業界・課題に合わせて選択。\nゼロから考える必要なし。", BLUE),
    ("圧倒的な低コスト", "個別開発の1/10以下のコストで\n導入可能。外部AI SaaS の\n月額コストも削減。", GREEN),
    ("即日〜1週間で稼働", "テンプレートベースだから\nPoCまで最短1日。\n本番稼働も1〜2週間。", ACCENT),
]

for i, (title, desc, color) in enumerate(values):
    x = Inches(0.8 + i * 4.0)
    bar = add_shape(slide, x, Inches(4.5), Inches(3.6), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.2), Inches(4.65), Inches(3.2), Inches(0.5),
                title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_textbox(slide, x + Inches(0.2), Inches(5.2), Inches(3.2), Inches(1.5),
                          [desc], font_size=13, color=DARK_GRAY, line_spacing=1.6,
                          alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 4: 全メニュー概要（21コース一覧）
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "メニュー全体構成：21コース × 170品",
             "どの業界のどの部門にも「すぐ出せる」ソリューションを網羅")
slide_footer(slide, 4)

table_data = [
    ["#", "コース名", "対象業界", "品数", "番号"],
    ["1", "共通コース", "全業種共通", "25品", "AIS-01〜25"],
    ["2", "製造業コース", "製造業", "8品", "AIS-26〜33"],
    ["3", "金融・保険コース", "銀行・証券・保険", "8品", "AIS-34〜41"],
    ["4", "医療・ヘルスケアコース", "病院・クリニック・薬局", "8品", "AIS-42〜49"],
    ["5", "小売・ECコース", "小売・EC事業者", "8品", "AIS-50〜57"],
    ["6", "不動産コース", "不動産会社・管理会社", "7品", "AIS-58〜64"],
    ["7", "教育コース", "学校・塾・研修機関", "7品", "AIS-65〜71"],
    ["8", "物流・SCMコース", "物流会社・メーカー物流部門", "6品", "AIS-72〜77"],
    ["9", "IT・SaaSコース", "IT企業・SaaS事業者", "7品", "AIS-78〜84"],
    ["10", "飲食・ホスピタリティコース", "飲食店・ホテル", "6品", "AIS-85〜90"],
    ["11", "建設コース", "建設会社・工務店", "6品", "AIS-91〜96"],
    ["12", "メディア・広告コース", "メディア・広告代理店", "6品", "AIS-97〜102"],
    ["13", "士業・専門サービスコース", "法律・会計・特許事務所", "6品", "AIS-103〜108"],
    ["14", "自治体・公共コース", "自治体・公共機関", "6品", "AIS-109〜114"],
    ["15", "人材・派遣コース", "人材紹介・派遣会社", "6品", "AIS-115〜120"],
]
# New courses highlighted
new_courses = [
    ["16", "クロスファンクション上級", "全業種（高度AI活用）", "10品", "AIS-121〜130"],
    ["17", "自動車・モビリティコース", "自動車メーカー・ディーラー", "8品", "AIS-131〜138"],
    ["18", "医薬品・ライフサイエンス", "製薬・CRO・医療機器", "8品", "AIS-139〜146"],
    ["19", "エネルギー・インフラ", "電力・ガス・再エネ", "8品", "AIS-147〜154"],
    ["20", "旅行・観光・ホテル", "旅行代理店・宿泊施設・DMO", "8品", "AIS-155〜162"],
    ["21", "農業・食品加工コース", "農業法人・食品メーカー", "8品", "AIS-163〜170"],
]
all_data = table_data + new_courses
all_data.append(["", "合計", "21コース", "170品", "AIS-01〜170"])

tbl = make_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), all_data,
                 [0.5, 3.2, 3.5, 0.8, 2.0], font_size=10, row_height=0.24)

# Highlight new courses (indices 16-21 in 0-based)
table_obj = tbl.table
total_rows = len(all_data)
for r_idx in range(16, total_rows - 1):
    for c_idx in range(5):
        cell = table_obj.cell(r_idx, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = ACCENT
            p.font.bold = True

# Total row (last row)
r_idx = total_rows - 1
for c_idx in range(5):
    cell = table_obj.cell(r_idx, c_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.bold = True

# Legend
add_rounded_rect(slide, Inches(0.5), Inches(7.0), Inches(4), Inches(0.35),
                 RGBColor(0xFF,0xF3,0xE0))
add_textbox(slide, Inches(0.7), Inches(7.0), Inches(3.5), Inches(0.35),
            "■ オレンジ = 今回の新規追加提案 50品", font_size=11, color=ACCENT, bold=True)


# ════════════════════════════════════════════════════
# SLIDE 5: 既存120品ハイライト（業界マップ）
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "既存120品：業界別 人気メニュー TOP3",
             "各業界で最もROIが高い「まず試すべき」ワークフロー")
slide_footer(slide, 5)

industry_picks = [
    ("共通", ["議事録AI要約", "社内ナレッジ検索", "営業提案書"], "全業種"),
    ("製造業", ["品質検査レポート", "SOP自動生成", "設備保全計画"], "品質・安全"),
    ("金融・保険", ["融資審査コメント", "AML/KYCチェック", "リスク管理"], "規制対応"),
    ("医療", ["診療サマリー", "患者説明文書", "医療安全レポート"], "文書作成"),
    ("小売・EC", ["商品説明文", "レビュー分析", "在庫分析"], "売上直結"),
    ("不動産", ["物件紹介文", "重説チェックリスト", "物件査定"], "業務効率"),
    ("教育", ["通知表所見文", "テスト問題生成", "授業計画"], "教員支援"),
    ("物流", ["配送ルート最適化", "通関書類チェック", "在庫最適化"], "コスト削減"),
    ("IT・SaaS", ["障害報告書", "テストケース生成", "コードレビュー"], "開発効率"),
    ("飲食", ["メニュー説明文", "衛生管理チェック", "原価分析"], "品質・利益"),
    ("建設", ["施工計画書", "安全管理計画", "工事見積"], "書類DX"),
    ("メディア", ["プレスリリース", "動画台本", "広告効果分析"], "コンテンツ"),
    ("士業", ["法律相談ドラフト", "経営コンサルレポート", "監査チェック"], "専門効率"),
    ("自治体", ["住民広報文", "議会答弁ドラフト", "補助金申請"], "住民SV"),
    ("人材", ["スカウトメール", "求人原稿一括", "人材紹介レポート"], "営業効率"),
]

cols = 5
rows_grid = 3
card_w = Inches(2.25)
card_h = Inches(1.65)

for idx, (ind, picks, tag) in enumerate(industry_picks):
    col = idx % cols
    row = idx // cols
    x = Inches(0.6 + col * 2.5)
    y = Inches(1.5 + row * 1.85)

    card = add_rounded_rect(slide, x, y, card_w, card_h, LIGHT_BLUE)
    # Industry name
    add_textbox(slide, x + Inches(0.1), y + Inches(0.05), Inches(1.5), Inches(0.3),
                ind, font_size=13, color=BLUE, bold=True)
    # Tag
    tag_box = add_rounded_rect(slide, x + Inches(1.55), y + Inches(0.05), Inches(0.65), Inches(0.25),
                               BLUE)
    add_textbox(slide, x + Inches(1.55), y + Inches(0.05), Inches(0.65), Inches(0.25),
                tag, font_size=8, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Picks
    for j, pick in enumerate(picks):
        add_textbox(slide, x + Inches(0.15), y + Inches(0.4 + j * 0.38), Inches(2.0), Inches(0.3),
                    f"  {pick}", font_size=10, color=DARK_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 6: 新規提案①クロスファンクション上級
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "【NEW】クロスファンクション上級コース 10品",
             "既存の単独AIサービスを Dify で再現し、外部コストを削減")
slide_footer(slide, 6)

cross_func_data = [
    ["#", "メニュー名", "再現対象（既存AIサービス）", "主な効果"],
    ["121", "AI文書校正・品質改善", "Grammarly / 文賢", "外部ツール月額×人数分を削減"],
    ["122", "商談録音AI分析・コーチング", "Gong / MiiTel", "営業1人月2h振り返り効率化"],
    ["123", "多言語翻訳・ローカライズ", "DeepL / WOVN", "翻訳外注コスト70%削減"],
    ["124", "請求書AI読取→仕訳提案", "freee / invox", "経理1人月10h入力作業削減"],
    ["125", "SNSブランドモニタリング", "Meltwater / Brandwatch", "月額10〜30万のツール代替"],
    ["126", "リードスコアリング", "Salesforce Einstein", "営業の無駄アプローチ50%削減"],
    ["127", "RFP自動回答ジェネレーター", "RFPIO / Loopio", "提案書作成時間70%短縮"],
    ["128", "ESGレポート生成", "Workiva / Diligent", "外部コンサル費用を大幅削減"],
    ["129", "経営ダッシュボードナラティブ", "Narrative Science", "月次報告書作成80%短縮"],
    ["130", "社内規程AI検索・改訂ドラフト", "（自社開発が必要だった領域）", "規程管理の属人化を解消"],
]
make_table(slide, Inches(0.4), Inches(1.5), Inches(12.5), cross_func_data,
           [0.5, 3.0, 3.5, 3.5], font_size=12, row_height=0.4)

# Callout
add_rounded_rect(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.65), RGBColor(0xFF,0xF3,0xE0))
add_textbox(slide, Inches(0.7), Inches(6.25), Inches(11.8), Inches(0.55),
            "Point：外部AI SaaSサービスの月額合計が1人あたり5,000〜20,000円 → Difyで一括代替すれば、100名規模で年間600万〜2,400万円のコスト削減",
            font_size=13, color=ACCENT, bold=True)


# ════════════════════════════════════════════════════
# SLIDE 7: 新規提案②自動車・モビリティコース
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "【NEW】自動車・モビリティコース 8品",
             "日本の基幹産業。品質文書・ディーラー業務をAIで革新")
slide_footer(slide, 7)

auto_data = [
    ["#", "メニュー名", "説明", "対象ユーザー"],
    ["131", "品質不具合分析（DRBFM）", "不具合→DRBFM形式の変化点分析・対策レポート自動生成", "品質管理部門"],
    ["132", "リコール影響範囲分析", "リコール情報→影響範囲・優先度・顧客通知文を生成", "品質保証・CS"],
    ["133", "FMEA分析シート生成", "設計/工程情報→故障モード・影響・検出のFMEAシート作成", "設計・生産技術"],
    ["134", "ディーラー接客スクリプト", "車種・顧客属性→試乗誘導・クロージングのスクリプト", "ディーラー営業"],
    ["135", "車両点検レポート生成", "点検データ→顧客向けレポート＋整備提案を自動作成", "サービスフロント"],
    ["136", "自動車保険見積説明書", "車両・顧客情報→保険見積の比較説明書を自動生成", "保険窓口"],
    ["137", "サービスキャンペーン通知", "キャンペーン情報→DM・メール・SMS通知文を一括生成", "マーケティング"],
    ["138", "IATF 16949監査チェック", "監査範囲→IATF要求事項の監査チェックリスト生成", "品質管理・監査"],
]
make_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), auto_data,
           [0.5, 2.8, 5.4, 1.8], font_size=11, row_height=0.55)


# ════════════════════════════════════════════════════
# SLIDE 8: 新規提案③医薬品・ライフサイエンスコース
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "【NEW】医薬品・ライフサイエンスコース 8品",
             "GxP文書・治験・薬事申請の膨大な文書作成をAIで効率化")
slide_footer(slide, 8)

pharma_data = [
    ["#", "メニュー名", "説明", "対象ユーザー"],
    ["139", "治験プロトコル要約生成", "治験計画書→IRB向け構造化要約（PICOT形式）を自動作成", "臨床開発部門"],
    ["140", "添付文書ドラフト生成", "薬理情報・臨床データ→添付文書ドラフトを自動生成", "薬事部門"],
    ["141", "副作用報告（CIOMS）", "副作用情報→CIOMS-I形式の安全性報告書ドラフト作成", "PV部門"],
    ["142", "MRディテーリング資料", "製品・競合情報→MR用ディテーリング資料を生成", "MR・営業"],
    ["143", "GMP逸脱報告書生成", "逸脱概要→GMP逸脱報告書（CAPA付き）を自動作成", "品質保証部門"],
    ["144", "安定性試験レポート", "試験データ→ICH Q1準拠の安定性試験レポートを生成", "品質管理部門"],
    ["145", "薬事申請CTD要約", "CTDデータ→品質・非臨床・臨床の概要ドラフト自動生成", "薬事申請担当"],
    ["146", "PVシグナル分析レポート", "有害事象データ→シグナル検出・評価・対応推奨を分析", "PV責任者"],
]
make_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), pharma_data,
           [0.5, 2.8, 5.4, 1.8], font_size=11, row_height=0.55)


# ════════════════════════════════════════════════════
# SLIDE 9: 新規提案④エネルギー・インフラコース
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "【NEW】エネルギー・インフラコース 8品",
             "脱炭素・電力自由化時代の安全管理・環境対応をAIで支援")
slide_footer(slide, 9)

energy_data = [
    ["#", "メニュー名", "説明", "対象ユーザー"],
    ["147", "プラント巡視点検レポート", "点検データ→異常検知ハイライト付きレポート自動作成", "保全担当"],
    ["148", "電力需給予測ナラティブ", "需給データ・気象予報→需給予測説明レポート自動生成", "需給運用"],
    ["149", "環境アセスメント報告書", "環境調査データ→環境影響評価報告書ドラフト生成", "環境部門"],
    ["150", "設備故障RCAレポート", "故障情報→根本原因分析＋再発防止策レポート自動作成", "設備保全"],
    ["151", "カーボンフットプリント計算", "排出データ→CO2排出量計算書（Scope1-3）を生成", "ESG推進"],
    ["152", "保安規程チェックリスト", "施設情報→電気/ガス事業法の保安規程チェック生成", "保安管理者"],
    ["153", "電力料金プラン提案書", "使用電力量→最適料金プラン比較＋提案書を自動作成", "営業担当"],
    ["154", "再エネ発電所月次報告", "発電データ→月次運転報告書を自動生成", "発電所管理"],
]
make_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), energy_data,
           [0.5, 2.8, 5.4, 1.8], font_size=11, row_height=0.55)


# ════════════════════════════════════════════════════
# SLIDE 10: 新規提案⑤旅行・観光・ホテルコース
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "【NEW】旅行・観光・ホテルコース 8品",
             "インバウンド4,000万人時代。多言語対応と顧客体験をAIで強化")
slide_footer(slide, 10)

travel_data = [
    ["#", "メニュー名", "説明", "対象ユーザー"],
    ["155", "旅行プラン自動提案", "希望条件→パーソナライズ旅行プラン＋見積を自動生成", "旅行カウンター"],
    ["156", "観光地多言語ガイド生成", "スポット情報→日英中韓の多言語ガイド文を一括生成", "DMO・観光協会"],
    ["157", "OTA口コミ分析・改善提案", "Booking/じゃらん口コミ→分析・改善優先度・回答文生成", "宿泊施設GM"],
    ["158", "宿泊プラン企画書", "施設特徴・ターゲット→プラン企画書（料金・販促文含む）", "企画担当"],
    ["159", "インバウンド接客フレーズ集", "業態・FAQ→英中韓の接客フレーズ＋対応マニュアル", "フロント・接客"],
    ["160", "ツアーガイド台本生成", "ルート情報→ガイドスクリプト（歴史・豆知識含む）作成", "ツアーガイド"],
    ["161", "旅行業約款チェックリスト", "ツアー内容→旅行業法に基づく説明事項チェック生成", "旅行業務取扱"],
    ["162", "観光マーケティングレポート", "観光客データ→観光DMP分析レポートを自動生成", "DMO・行政"],
]
make_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), travel_data,
           [0.5, 2.8, 5.4, 1.8], font_size=11, row_height=0.55)


# ════════════════════════════════════════════════════
# SLIDE 11: 新規提案⑥農業・食品加工コース
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "【NEW】農業・食品加工コース 8品",
             "スマート農業・HACCP対応・6次産業化をAIで推進")
slide_footer(slide, 11)

agri_data = [
    ["#", "メニュー名", "説明", "対象ユーザー"],
    ["163", "栽培記録・出荷記録生成", "作物情報・農薬使用→GAP準拠の栽培・出荷記録作成", "農業法人"],
    ["164", "HACCP管理記録シート", "加工工程・CCP情報→HACCP対応の管理記録シート生成", "食品工場"],
    ["165", "農産物POP・商品説明文", "農産物特徴→直売所POP＋EC用商品説明を自動生成", "直売所・EC担当"],
    ["166", "食品表示ラベルチェック", "原材料・アレルゲン→食品表示法準拠のチェック結果生成", "品質管理"],
    ["167", "農水省補助金申請ドラフト", "事業計画→農水省系補助金の申請書ドラフト自動作成", "経営企画"],
    ["168", "病害虫診断・対策提案", "症状・栽培環境→病害虫の推定診断＋防除対策を提案", "営農指導員"],
    ["169", "トレーサビリティ報告書", "生産〜流通データ→食品トレーサビリティ報告書を生成", "品質保証"],
    ["170", "6次産業化事業計画書", "農産物・加工品→6次産業化の事業計画書ドラフト作成", "農業経営者"],
]
make_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), agri_data,
           [0.5, 2.8, 5.4, 1.8], font_size=11, row_height=0.55)


# ════════════════════════════════════════════════════
# SLIDE 12: 導入効果・ROI
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "導入効果シミュレーション",
             "ワークフロー1品あたりの平均ROIは年間900万円")
slide_footer(slide, 12)

# Left: ROI calculation
add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(3.2), LIGHT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.5),
            "1品あたりの年間コスト削減効果（試算）", font_size=18, color=BLUE, bold=True)

roi_lines = [
    "対象担当者数:                    10名",
    "利用頻度:                              週2回",
    "1回あたり時間削減:            1.8時間（3h → 1.2h）",
    "年間削減時間:                      1,800時間",
    "時給換算（@5,000円）:       900万円 / 年",
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2.2),
                      [(l, 14, DARK_GRAY, False) for l in roi_lines],
                      line_spacing=1.8)

# Right: Package ROI
add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5),
            "導入パターン別 年間効果", font_size=18, color=NAVY, bold=True)

pkg_data = [
    ["導入パターン", "品数", "コスト削減", "売上向上", "年間合計"],
    ["スモールスタート", "5品", "600万", "—", "600万"],
    ["部門導入", "15品", "1,500万", "300万", "1,800万"],
    ["全社導入", "25品", "3,000万", "1,000万", "4,000万"],
    ["フル活用", "30品+", "5,000万", "2,000万", "7,000万"],
]
make_table(slide, Inches(7.0), Inches(2.2), Inches(5.8), pkg_data,
           [1.5, 0.7, 1.0, 1.0, 1.0], font_size=12, row_height=0.45)

# Bottom: Key message
add_rounded_rect(slide, Inches(0.6), Inches(5.0), Inches(12.0), Inches(1.8), RGBColor(0xE8,0xF5,0xE9))
add_textbox(slide, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.5),
            "外部AI SaaSサービスの代替効果（クロスファンクション上級コース導入時）",
            font_size=16, color=GREEN, bold=True)
add_multiline_textbox(slide, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0), [
    ("Grammarly Business: 月額$25/人 × 100名 = 年間約450万円 → Difyで代替", 13, DARK_GRAY, False),
    ("翻訳外注: 月50万円 × 12ヶ月 = 年間600万円 → 70%削減で420万円の効果", 13, DARK_GRAY, False),
    ("ブランドモニタリングツール: 月額30万 × 12ヶ月 = 年間360万円 → Difyで代替", 13, DARK_GRAY, False),
    ("合計: 年間1,200万円以上の外部サービスコストを Dify に一本化可能", 14, GREEN, True),
], line_spacing=1.5)


# ════════════════════════════════════════════════════
# SLIDE 13: 選び方ガイド（フローチャート風）
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "選び方ガイド：業界×課題マトリクス",
             "「何から始めるべきか」を一目で判断")
slide_footer(slide, 13)

matrix_data = [
    ["業界", "まず試すべき1品", "次に広げる2品目", "最終ゴール"],
    ["製造業", "AIS-26 品質検査レポート", "AIS-27 SOP自動生成", "品質文書の完全自動化"],
    ["金融・保険", "AIS-34 融資審査コメント", "AIS-37 AML/KYCチェック", "規制対応のスピード化"],
    ["医療", "AIS-43 診療サマリー", "AIS-42 患者説明文書", "医師の文書作成ゼロ"],
    ["小売・EC", "AIS-50 商品説明文", "AIS-51 レビュー分析", "データドリブン経営"],
    ["不動産", "AIS-58 物件紹介文", "AIS-59 重説チェック", "物件登録スピード3倍"],
    ["教育", "AIS-68 通知表所見文", "AIS-66 テスト問題生成", "教員の働き方改革"],
    ["IT・SaaS", "AIS-78 障害報告書", "AIS-81 テストケース", "ドキュメント負債の解消"],
    ["自動車【NEW】", "AIS-131 DRBFM分析", "AIS-133 FMEA生成", "品質管理DX"],
    ["医薬品【NEW】", "AIS-139 治験プロトコル", "AIS-143 GMP逸脱報告", "GxP文書の効率化"],
    ["エネルギー【NEW】", "AIS-147 巡視点検レポート", "AIS-151 CO2計算書", "脱炭素対応の加速"],
    ["旅行・観光【NEW】", "AIS-155 旅行プラン提案", "AIS-156 多言語ガイド", "インバウンド対応完備"],
    ["農業【NEW】", "AIS-163 栽培記録", "AIS-164 HACCP記録", "スマート農業の実現"],
]
make_table(slide, Inches(0.4), Inches(1.5), Inches(12.5), matrix_data,
           [1.8, 3.0, 3.0, 2.8], font_size=11, row_height=0.38)


# ════════════════════════════════════════════════════
# SLIDE 14: 導入ステップ
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "導入ロードマップ：5ステップ",
             "PoC（概念実証）から全社展開まで、段階的に進めます")
slide_footer(slide, 14)

steps = [
    ("Step 1", "課題特定", "1日", "現場ヒアリングで\n対象業務を選定", BLUE),
    ("Step 2", "PoC構築", "1〜2週間", "メニューから2〜3品を\n選びプロトタイプ構築", RGBColor(0x00,0x96,0x88)),
    ("Step 3", "パイロット", "1ヶ月", "限定チーム(5-10名)で\n実運用・KPI測定", GREEN),
    ("Step 4", "本格導入", "2〜3ヶ月", "全社展開\n追加ワークフロー導入", GOLD),
    ("Step 5", "拡張最適化", "継続", "プロンプトチューニング\n他部門・他拠点へ横展開", ACCENT),
]

for i, (step_name, title, period, desc, color) in enumerate(steps):
    x = Inches(0.4 + i * 2.5)
    y = Inches(1.8)

    # Step circle
    add_circle_number(slide, x + Inches(0.7), y, Inches(0.8), str(i+1), color)

    # Step name
    add_textbox(slide, x, y + Inches(0.95), Inches(2.2), Inches(0.4),
                title, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Period badge
    badge = add_rounded_rect(slide, x + Inches(0.4), y + Inches(1.4), Inches(1.4), Inches(0.3), color)
    add_textbox(slide, x + Inches(0.4), y + Inches(1.4), Inches(1.4), Inches(0.3),
                period, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_multiline_textbox(slide, x + Inches(0.1), y + Inches(1.85), Inches(2.0), Inches(1.2),
                          [desc], font_size=12, color=DARK_GRAY, line_spacing=1.5,
                          alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 4:
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, x + Inches(2.2), y + Inches(0.25),
            Inches(0.3), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xCC,0xCC,0xCC)
        arrow.line.fill.background()

# Bottom KPI boxes
add_shape(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.04), BLUE)
add_textbox(slide, Inches(0.4), Inches(5.15), Inches(12), Inches(0.4),
            "パイロット運用の成功基準（KPI）", font_size=16, color=NAVY, bold=True)

kpis = [
    ("時間削減", "対象業務の\n作成時間 50%↓"),
    ("品質", "出力修正率\n30%以下"),
    ("利用率", "対象者の\n80%以上が利用"),
    ("満足度", "ユーザー満足度\n4.0 / 5.0 以上"),
]
for i, (kpi_name, kpi_desc) in enumerate(kpis):
    x = Inches(0.6 + i * 3.1)
    card = add_rounded_rect(slide, x, Inches(5.6), Inches(2.7), Inches(1.3), LIGHT_BLUE)
    add_textbox(slide, x + Inches(0.2), Inches(5.7), Inches(2.3), Inches(0.35),
                kpi_name, font_size=14, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_textbox(slide, x + Inches(0.2), Inches(6.05), Inches(2.3), Inches(0.7),
                          [kpi_desc], font_size=13, color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                          line_spacing=1.4)


# ════════════════════════════════════════════════════
# SLIDE 15: まとめ・Next Steps
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT)
add_shape(slide, Inches(0), Inches(7.42), W, Inches(0.08), ACCENT)

add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(0.8),
            "まとめ", font_size=36, color=WHITE, bold=True)
add_shape(slide, Inches(1.5), Inches(1.5), Inches(2), Inches(0.04), ACCENT)

# Summary stats - 3 big numbers
stats = [
    ("170", "品", "AIワークフロー"),
    ("21", "コース", "業界カバー"),
    ("7,000", "万円", "年間最大ROI"),
]
for i, (num, unit, label) in enumerate(stats):
    x = Inches(1.2 + i * 3.8)
    add_textbox(slide, x, Inches(2.0), Inches(3.2), Inches(1.0),
                num, font_size=60, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(2.9), Inches(3.2), Inches(0.4),
                unit, font_size=20, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, Inches(3.3), Inches(3.2), Inches(0.4),
                label, font_size=16, color=RGBColor(0xAA,0xCC,0xFF), alignment=PP_ALIGN.CENTER)

# Key messages
add_multiline_textbox(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(1.5), [
    ("「うちの業界の、この業務に、このAIワークフローが使える！」", 18, WHITE, True),
    ("", 6, WHITE, False),
    ("170品のファミレスメニューがあれば、どの業界のどの部門の顧客に対しても", 16, RGBColor(0xCC,0xDD,0xFF), False),
    ("「御社に最適なAIワークフローが、すでにご用意されています」と即座に提案できます。", 16, RGBColor(0xCC,0xDD,0xFF), False),
], line_spacing=1.5)

# Next Steps
add_shape(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.04), ACCENT)
add_textbox(slide, Inches(1.5), Inches(5.95), Inches(3), Inches(0.4),
            "Next Steps", font_size=20, color=ACCENT, bold=True)

add_multiline_textbox(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.9), [
    ("Step 1: メニューから2〜3品を選定 → Step 2: 1〜2週間でPoC構築 → Step 3: 効果を実感", 14, RGBColor(0xCC,0xDD,0xFF), False),
], line_spacing=1.5)


# ════════════════════════════════════════════════════
# 保存
# ════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Dify_AI_Workflow_FamilyRestaurant_Menu.pptx")
prs.save(output_path)
print(f"PowerPoint saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
