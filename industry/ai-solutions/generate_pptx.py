#!/usr/bin/env python3
"""
AI Solutions Catalog PowerPoint Generator
全120品のAIソリューションカタログをPowerPointで概観できるように生成
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# カラー定義
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
NAVY = RGBColor(0x1A, 0x23, 0x7E)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x19, 0x76, 0xD2)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)

def hex_to_rgb(hex_str):
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def darken(hex_str, factor=0.6):
    h = hex_str.lstrip('#')
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(int(r * factor), int(g * factor), int(b * factor))

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# データ定義
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
COURSES = [
    {
        "name": "共通コース（全業種共通）",
        "icon": "🏢",
        "bg": "#E3F2FD",
        "count": 25,
        "range": "AIS-01〜25",
        "desc": "全業種で必要な基盤業務を支えるAIワークフロー",
        "categories": [
            ("マーケティング&セールス", [
                (1, "SEO記事自動生成"),
                (8, "マーケティングコピー一括生成"),
                (9, "AI営業提案書ジェネレーター"),
                (11, "メール文面パーソナライズ生成"),
                (12, "競合分析レポート生成"),
                (13, "LP/広告コピーA/Bテスト生成"),
                (22, "見積書ドラフト生成"),
            ]),
            ("カスタマーサポート", [
                (3, "CS FAQチャットボット"),
                (14, "VOC分析・感情分析レポート"),
                (15, "問い合わせ自動分類・優先度判定"),
            ]),
            ("HR・人事", [
                (10, "AI求人票・JD自動生成"),
                (18, "履歴書スクリーニング・評価"),
                (19, "面接質問自動生成"),
                (20, "人事評価コメント生成"),
                (25, "社内研修コンテンツ生成"),
            ]),
            ("法務・コンプライアンス", [
                (2, "契約書リスクレビュー"),
                (16, "利用規約・プライバシーポリシー生成"),
                (17, "コンプライアンスチェックリスト生成"),
            ]),
            ("バックオフィス・経営", [
                (4, "議事録AI要約・アクション抽出"),
                (6, "長文ドキュメントAI要約"),
                (7, "WebリサーチAIレポート自動生成"),
                (21, "財務分析レポート生成"),
            ]),
            ("ナレッジ・IT", [
                (5, "社内ナレッジAI検索チャット"),
                (23, "Text-to-SQL クエリ生成"),
                (24, "API仕様書自動生成"),
            ]),
        ],
    },
    {
        "name": "製造業コース",
        "icon": "🏭",
        "bg": "#E0F7FA",
        "count": 8,
        "range": "AIS-26〜33",
        "desc": "品質管理・工程管理・安全衛生の文書作成を効率化",
        "items": [
            (26, "品質検査レポート自動生成"),
            (27, "作業手順書(SOP)AI生成"),
            (28, "設備保全計画レポート"),
            (29, "サプライヤー評価シート生成"),
            (30, "製品仕様書ドラフト生成"),
            (31, "生産計画アドバイザー"),
            (32, "製造クレーム分析レポート"),
            (33, "安全衛生パトロールチェックリスト"),
        ],
    },
    {
        "name": "金融・保険コース",
        "icon": "💰",
        "bg": "#FFF8E1",
        "count": 8,
        "range": "AIS-34〜41",
        "desc": "審査・コンプライアンス・顧客対応の正確性と迅速性",
        "items": [
            (34, "融資審査コメント生成"),
            (35, "金融商品説明書AI生成"),
            (36, "保険商品レコメンドレポート"),
            (37, "AML/KYCチェックレポート"),
            (38, "投資分析レポート生成"),
            (39, "顧客ポートフォリオ分析"),
            (40, "保険クレーム査定コメント"),
            (41, "リスク管理レポート生成"),
        ],
    },
    {
        "name": "医療・ヘルスケアコース",
        "icon": "🏥",
        "bg": "#E8F5E9",
        "count": 8,
        "range": "AIS-42〜49",
        "desc": "患者のわかりやすさと医療者の効率を両立",
        "items": [
            (42, "患者説明文書生成"),
            (43, "診療サマリーAI生成"),
            (44, "医療論文サマリー"),
            (45, "薬剤情報提供書生成"),
            (46, "AI問診票分析"),
            (47, "医療安全インシデントレポート"),
            (48, "栄養指導計画AI生成"),
            (49, "リハビリテーション計画書生成"),
        ],
    },
    {
        "name": "小売・ECコース",
        "icon": "🛒",
        "bg": "#FCE4EC",
        "count": 8,
        "range": "AIS-50〜57",
        "desc": "商品の魅力を伝え、顧客を理解し、在庫を最適化",
        "items": [
            (50, "商品説明文AI生成"),
            (51, "レビュー分析・改善提案"),
            (52, "在庫分析・発注推奨レポート"),
            (53, "プライシング分析レポート"),
            (54, "顧客セグメント分析"),
            (55, "メルマガ・LINE配信文生成"),
            (56, "商品FAQ自動生成"),
            (57, "接客マニュアルAI生成"),
        ],
    },
    {
        "name": "不動産コース",
        "icon": "🏠",
        "bg": "#E8F5E9",
        "count": 7,
        "range": "AIS-58〜64",
        "desc": "物件情報の見せ方と取引の抜け漏れ防止を同時に",
        "items": [
            (58, "物件紹介文自動生成"),
            (59, "重要事項説明書チェックリスト"),
            (60, "物件査定レポートAI生成"),
            (61, "内覧フォローメール生成"),
            (62, "賃貸契約書チェックリスト"),
            (63, "不動産投資分析レポート"),
            (64, "管理物件月次レポート"),
        ],
    },
    {
        "name": "教育コース",
        "icon": "📚",
        "bg": "#FFF9C4",
        "count": 7,
        "range": "AIS-65〜71",
        "desc": "教育の質を高めながら先生の働き方改革を実現",
        "items": [
            (65, "授業計画（レッスンプラン）AI生成"),
            (66, "テスト問題自動生成"),
            (67, "学習進捗分析レポート"),
            (68, "通知表所見文AI生成"),
            (69, "シラバス自動生成"),
            (70, "学習教材レコメンド"),
            (71, "保護者通知文生成"),
        ],
    },
    {
        "name": "物流・サプライチェーンコース",
        "icon": "🚛",
        "bg": "#E3F2FD",
        "count": 6,
        "range": "AIS-72〜77",
        "desc": "物流の見える化と標準化をAIで加速",
        "items": [
            (72, "配送ルート最適化提案"),
            (73, "倉庫作業手順書AI生成"),
            (74, "通関書類チェックリスト"),
            (75, "サプライチェーンリスク分析"),
            (76, "輸送コスト分析レポート"),
            (77, "在庫最適化AI提案"),
        ],
    },
    {
        "name": "IT・SaaSコース",
        "icon": "💻",
        "bg": "#E8EAF6",
        "count": 7,
        "range": "AIS-78〜84",
        "desc": "開発チームのドキュメント負債をAIで返済",
        "items": [
            (78, "障害報告書（インシデントレポート）生成"),
            (79, "リリースノート自動生成"),
            (80, "コードレビューコメント生成"),
            (81, "テストケース自動生成"),
            (82, "セキュリティ診断チェックリスト"),
            (83, "SLA/SLOレポート自動生成"),
            (84, "ユーザーストーリー生成"),
        ],
    },
    {
        "name": "飲食・ホスピタリティコース",
        "icon": "🍽️",
        "bg": "#FFF3E0",
        "count": 6,
        "range": "AIS-85〜90",
        "desc": "美味しいを伝わる言葉に、安全を仕組みに",
        "items": [
            (85, "メニュー説明文AI生成"),
            (86, "口コミ返信文自動生成"),
            (87, "衛生管理チェックリスト生成"),
            (88, "宴会・コースプラン提案"),
            (89, "スタッフシフト最適化提案"),
            (90, "食材原価分析レポート"),
        ],
    },
    {
        "name": "建設コース",
        "icon": "🏗️",
        "bg": "#EFEBE9",
        "count": 6,
        "range": "AIS-91〜96",
        "desc": "施工現場の安全と品質を文書で守る",
        "items": [
            (91, "施工計画書ドラフト生成"),
            (92, "安全管理計画書AI生成"),
            (93, "工事見積明細AI生成"),
            (94, "施工品質チェックリスト"),
            (95, "建設プロジェクト週報生成"),
            (96, "近隣説明文書AI生成"),
        ],
    },
    {
        "name": "メディア・広告コース",
        "icon": "📺",
        "bg": "#F3E5F5",
        "count": 6,
        "range": "AIS-97〜102",
        "desc": "コンテンツ制作の量と質を同時に引き上げる",
        "items": [
            (97, "プレスリリース自動生成"),
            (98, "SNS投稿カレンダー生成"),
            (99, "インフルエンサー分析レポート"),
            (100, "動画台本（スクリプト）生成"),
            (101, "ポッドキャスト台本生成"),
            (102, "広告効果分析レポート"),
        ],
    },
    {
        "name": "士業・専門サービスコース",
        "icon": "⚖️",
        "bg": "#FBE9E7",
        "count": 6,
        "range": "AIS-103〜108",
        "desc": "高度な専門知識を活かし定型業務のドラフトをAI支援",
        "items": [
            (103, "法律相談回答ドラフト"),
            (104, "税務相談回答テンプレート"),
            (105, "特許出願明細書ドラフト"),
            (106, "会計監査チェックリスト"),
            (107, "経営コンサルレポート生成"),
            (108, "顧問契約提案書生成"),
        ],
    },
    {
        "name": "自治体・公共コース",
        "icon": "🏛️",
        "bg": "#E0F2F1",
        "count": 6,
        "range": "AIS-109〜114",
        "desc": "住民サービス向上と職員の業務効率化を両立",
        "items": [
            (109, "住民向け広報文生成"),
            (110, "政策ブリーフィング資料生成"),
            (111, "補助金申請書ドラフト"),
            (112, "公共施設FAQ生成"),
            (113, "防災マニュアルAI生成"),
            (114, "議会答弁ドラフト生成"),
        ],
    },
    {
        "name": "人材・派遣コース",
        "icon": "👥",
        "bg": "#FCE4EC",
        "count": 6,
        "range": "AIS-115〜120",
        "desc": "求人・スカウト・推薦のコミュニケーションをAI強化",
        "items": [
            (115, "求人原稿一括生成"),
            (116, "スカウトメール文面生成"),
            (117, "候補者サーチレポート"),
            (118, "派遣契約書チェックリスト"),
            (119, "人材紹介レポート生成"),
            (120, "キャリアカウンセリング提案"),
        ],
    },
]


def add_bg_rect(slide, color):
    """スライド全体の背景色を設定"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, line_color=None):
    """角丸四角形を追加"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def set_text(shape, text, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    """シェイプのテキスト設定"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf


def add_text_para(tf, text, font_size=12, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, space_before=0):
    """テキストフレームに段落を追加"""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    return p


def create_title_slide(prs):
    """タイトルスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg_rect(slide, NAVY)

    # メインタイトル
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI Solutions Catalog"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_text_para(tf, "業界別 AI ワークフロー カタログ 全120品", 24, False, RGBColor(0xBB, 0xDE, 0xFB), PP_ALIGN.CENTER, 12)
    add_text_para(tf, "Powered by Dify", 18, False, RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER, 24)

    # サブテキスト
    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(7), Inches(1))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "「うちの業界の、この業務に、このAIワークフローが使える！」"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    p2.alignment = PP_ALIGN.CENTER

    add_text_para(tf2, 'そんな "あ、これだ！" の発見を提供する、業界特化型AIソリューションメニュー', 14, False, RGBColor(0xBB, 0xDE, 0xFB), PP_ALIGN.CENTER, 8)


def create_overview_slide(prs):
    """全体概要スライド - 15コースの一覧"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, WHITE)

    # タイトル
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "全15コース・120品 メニュー概要"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT

    # 5列 x 3行のグリッドでカードを配置
    cols = 5
    card_w = Inches(1.8)
    card_h = Inches(1.55)
    margin_x = Inches(0.15)
    margin_y = Inches(0.12)
    start_x = Inches(0.25)
    start_y = Inches(1.1)

    for i, course in enumerate(COURSES):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + margin_x)
        y = start_y + row * (card_h + margin_y)

        bg_color = hex_to_rgb(course["bg"])
        card = add_shape(slide, x, y, card_w, card_h, bg_color)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(8)
        tf.margin_right = Pt(8)
        tf.margin_top = Pt(6)
        tf.margin_bottom = Pt(4)

        # アイコン + コース名
        p = tf.paragraphs[0]
        p.text = f"{course['icon']} {course['name'].replace('コース', '').strip()}"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

        # 品数
        add_text_para(tf, f"{course['count']}品 | {course['range']}", 7, False, MED_GRAY, PP_ALIGN.CENTER, 2)

        # 説明
        add_text_para(tf, course['desc'], 6.5, False, MED_GRAY, PP_ALIGN.CENTER, 3)


def create_course_slide(prs, course):
    """各コースの詳細スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, WHITE)

    bg_color = hex_to_rgb(course["bg"])
    title_color = darken(course["bg"], 0.4)

    # ヘッダーバー
    header = add_shape(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), bg_color)
    tf = header.text_frame
    tf.margin_left = Pt(20)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = f"{course['icon']}  {course['name']}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = title_color
    add_text_para(tf, f"{course['desc']}  |  {course['count']}品  |  {course['range']}", 11, False, MED_GRAY, PP_ALIGN.LEFT, 2)

    if "categories" in course:
        # 共通コースはカテゴリ別に表示
        y_pos = Inches(1.05)
        for cat_name, items in course["categories"]:
            # カテゴリ名
            txBox = slide.shapes.add_textbox(Inches(0.4), y_pos, Inches(9), Inches(0.28))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"■ {cat_name}"
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
            y_pos += Inches(0.28)

            # アイテム (横並び)
            items_per_row = 4
            for j in range(0, len(items), items_per_row):
                row_items = items[j:j+items_per_row]
                for k, (num, name) in enumerate(row_items):
                    x = Inches(0.4) + k * Inches(2.3)
                    card = add_shape(slide, x, y_pos, Inches(2.2), Inches(0.35), LIGHT_GRAY)
                    tf2 = card.text_frame
                    tf2.margin_left = Pt(6)
                    tf2.margin_top = Pt(2)
                    tf2.word_wrap = True
                    p2 = tf2.paragraphs[0]
                    p2.text = f"AIS-{num:02d}"
                    p2.font.size = Pt(7)
                    p2.font.bold = True
                    p2.font.color.rgb = ACCENT_BLUE
                    add_text_para(tf2, name, 7.5, False, DARK_GRAY, PP_ALIGN.LEFT, 0)
                y_pos += Inches(0.4)
            y_pos += Inches(0.05)
    else:
        # 業界コース - カード形式
        items = course["items"]
        cols = 2
        card_w = Inches(4.5)
        card_h = Inches(0.55)
        margin_x = Inches(0.2)
        margin_y = Inches(0.1)
        start_x = Inches(0.4)
        start_y = Inches(1.1)

        for j, (num, name) in enumerate(items):
            row = j // cols
            col = j % cols
            x = start_x + col * (card_w + margin_x)
            y = start_y + row * (card_h + margin_y)

            card = add_shape(slide, x, y, card_w, card_h, LIGHT_GRAY)
            tf2 = card.text_frame
            tf2.word_wrap = True
            tf2.margin_left = Pt(12)
            tf2.margin_top = Pt(6)

            p2 = tf2.paragraphs[0]
            p2.text = f"AIS-{num:02d}"
            p2.font.size = Pt(10)
            p2.font.bold = True
            p2.font.color.rgb = title_color

            add_text_para(tf2, name, 12, False, DARK_GRAY, PP_ALIGN.LEFT, 1)


def create_matrix_slide(prs):
    """業界 x 課題マトリクス スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, WHITE)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "業界別 おすすめスタートガイド"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    matrix = [
        ("🏭 製造業", "AIS-26 品質検査レポート", "AIS-27 SOP"),
        ("💰 金融・保険", "AIS-34 融資審査コメント", "AIS-37 AML/KYC"),
        ("🏥 医療", "AIS-43 診療サマリー", "AIS-42 患者説明文書"),
        ("🛒 小売・EC", "AIS-50 商品説明文", "AIS-51 レビュー分析"),
        ("🏠 不動産", "AIS-58 物件紹介文", "AIS-59 重説チェック"),
        ("📚 教育", "AIS-68 通知表所見文", "AIS-66 テスト問題"),
        ("🚛 物流", "AIS-72 配送ルート最適化", "AIS-74 通関チェック"),
        ("💻 IT・SaaS", "AIS-78 障害報告書", "AIS-81 テストケース"),
        ("🍽️ 飲食", "AIS-85 メニュー説明文", "AIS-87 衛生管理"),
        ("🏗️ 建設", "AIS-91 施工計画書", "AIS-92 安全管理計画"),
        ("📺 メディア", "AIS-97 プレスリリース", "AIS-100 動画台本"),
        ("⚖️ 士業", "AIS-103 法律相談回答", "AIS-107 コンサルレポート"),
        ("🏛️ 自治体", "AIS-109 住民向け広報文", "AIS-114 議会答弁"),
        ("👥 人材", "AIS-116 スカウトメール", "AIS-115 求人原稿"),
    ]

    # ヘッダー行
    headers = ["業界", "まず試すべき1品", "次に広げる2品目"]
    header_widths = [Inches(1.8), Inches(3.5), Inches(3.5)]
    x_start = Inches(0.5)
    y_start = Inches(1.1)
    row_h = Inches(0.35)

    for k, (htext, hw) in enumerate(zip(headers, header_widths)):
        x = x_start + sum(w for w in [Inches(0)] + list(header_widths[:k]))
        shape = add_shape(slide, x, y_start, hw, row_h, NAVY)
        set_text(shape, htext, 10, True, WHITE, PP_ALIGN.CENTER)

    for i, (industry, first, second) in enumerate(matrix):
        y = y_start + row_h + i * Inches(0.32)
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        vals = [industry, first, second]
        for k, (val, hw) in enumerate(zip(vals, header_widths)):
            x = x_start + sum(w for w in [Inches(0)] + list(header_widths[:k]))
            shape = add_shape(slide, x, y, hw, Inches(0.30), bg)
            fs = 9 if k == 0 else 8.5
            bold = k == 0
            set_text(shape, val, fs, bold, DARK_GRAY, PP_ALIGN.CENTER)


def create_steps_slide(prs):
    """導入ステップ スライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, WHITE)

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Dify AI ワークフロー 導入5ステップ"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    steps = [
        ("Step 1", "課題の特定", "1日", "現場ヒアリングで対象業務を選定\n週1回以上・30分以上の文書作成業務を優先"),
        ("Step 2", "PoC（概念実証）", "1-2週間", "カタログから2-3品を選びプロトタイプ構築\n実際の業務データで効果を検証"),
        ("Step 3", "パイロット運用", "1ヶ月", "5-10名の限定チームで実運用開始\nKPI: 時間50%短縮・修正率30%以下"),
        ("Step 4", "本格導入", "2-3ヶ月", "全社展開・トレーニング実施\n利用ガイドライン・テンプレート整備"),
        ("Step 5", "拡張・最適化", "継続", "プロンプトチューニング・品質向上\n新規ワークフロー追加・横展開"),
    ]

    for i, (step, title, period, desc) in enumerate(steps):
        x = Inches(0.3) + i * Inches(1.9)
        y = Inches(1.3)

        # ステップカード
        card = add_shape(slide, x, y, Inches(1.8), Inches(3.5), LIGHT_GRAY)
        tf2 = card.text_frame
        tf2.word_wrap = True
        tf2.margin_left = Pt(10)
        tf2.margin_right = Pt(10)
        tf2.margin_top = Pt(10)

        p2 = tf2.paragraphs[0]
        p2.text = step
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = ACCENT_BLUE
        p2.alignment = PP_ALIGN.CENTER

        add_text_para(tf2, title, 13, True, DARK_GRAY, PP_ALIGN.CENTER, 6)
        add_text_para(tf2, f"({period})", 9, False, MED_GRAY, PP_ALIGN.CENTER, 2)
        add_text_para(tf2, "", 6, False, MED_GRAY, PP_ALIGN.CENTER, 8)

        for line in desc.split('\n'):
            add_text_para(tf2, line, 8, False, MED_GRAY, PP_ALIGN.LEFT, 4)


def create_closing_slide(prs):
    """クロージングスライド"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, NAVY)

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "全120品のAIワークフローで"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_text_para(tf, "御社の業務を変革します", 32, True, WHITE, PP_ALIGN.CENTER, 8)

    txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "15業界 × ノーコード/ローコード × 最短1日でプロトタイプ完成"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
    p2.alignment = PP_ALIGN.CENTER

    add_text_para(tf2, "", 10, False, WHITE, PP_ALIGN.CENTER, 16)
    add_text_para(tf2, "Powered by Dify  |  AI Solutions Menu Catalog v2.0", 14, False, RGBColor(0x90, 0xCA, 0xF9), PP_ALIGN.CENTER, 8)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    # 1. タイトル
    create_title_slide(prs)

    # 2. 全体概要
    create_overview_slide(prs)

    # 3. 各コースの詳細
    for course in COURSES:
        create_course_slide(prs, course)

    # 4. おすすめマトリクス
    create_matrix_slide(prs)

    # 5. 導入ステップ
    create_steps_slide(prs)

    # 6. クロージング
    create_closing_slide(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AI_Solutions_Catalog_v2.pptx")
    prs.save(output_path)
    print(f"Generated: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
